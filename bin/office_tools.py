#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
office_tools.py — Office 文件处理（Excel 优先 + pandoc 写作）
科研骨架模块 | 依赖: openpyxl, pandoc(写作)

用法:
  python office_tools.py read data.xlsx [--sheet 表1] [--head 10]
  python office_tools.py md2xlsx 笔记.md out.xlsx [--table all|0|1...]
  python office_tools.py xlsx2md data.xlsx [--sheet 表1]
  python office_tools.py stats data.xlsx [--sheet 表1] [--col 列名,列名]
  python office_tools.py csv2xlsx data.csv out.xlsx
  python office_tools.py extract pdf 论文.pdf --outdir 图/ [--pages 1,3] [--min-size 250] [--min-kb 5]
  python office_tools.py extract docx 报告.docx --outdir 图/
  python office_tools.py extract pptx 汇报.pptx --outdir 图/
  python office_tools.py classify 图/        # vision 分类幸存图：价值图 vs 装饰图
  python office_tools.py md2docx 笔记.md out.docx [--toc] [--reference-doc 模板.docx]
  python office_tools.py md2pptx 讲稿.md out.pptx [--slide-level 2] [--reference-doc 模板.pptx]

extract 免费过滤: 尺寸/文件大小/页眉页脚 + 「Figure N」标题邻近标记([CAP]); 只对幸存图用 vision(classify)

md→docx/pptx 走 pandoc: LaTeX 公式($..$)→Word/PPT 原生 OMML 方程; --reference-doc 控样式(中文宋体/黑体)
场景: 阅读报告/防撞车矩阵(markdown 表) → Excel 分析；实验数据 → 统计；Excel → markdown 回写笔记；markdown 写作 → Word/PPT
"""

import argparse, csv, os, re, shutil, subprocess, sys
from openpyxl import Workbook, load_workbook

sys.stdout.reconfigure(encoding="utf-8")


# ---------- 读取 ----------

def list_sheets(path):
    return load_workbook(path, read_only=True).sheetnames


def read_sheet(path, sheet=None, head=0):
    wb = load_workbook(path, data_only=True)
    if sheet is None:
        sheet = wb.sheetnames[0]
    ws = wb[sheet]
    rows = [["" if v is None else str(v) for v in row] for row in ws.iter_rows(values_only=True)]
    if head > 0:
        rows = rows[:head]
    return sheet, rows


# ---------- markdown 表格解析 ----------

def extract_tables(md_text):
    """把 md 文本切成若干张连续 '|' 行组成的表。"""
    tables, cur = [], []
    for l in md_text.splitlines():
        if l.strip().startswith("|"):
            cur.append(l)
        elif cur:
            tables.append(cur)
            cur = []
    if cur:
        tables.append(cur)
    return tables


def md_table_to_rows(table_lines):
    rows = []
    for l in table_lines:
        cells = [c.strip() for c in l.strip().strip("|").split("|")]
        if all(re.fullmatch(r":?-{2,}:?", c) for c in cells):
            continue  # 分隔行 |---|
        rows.append(cells)
    return rows


# ---------- 命令 ----------

def cmd_read(args):
    sheet, rows = read_sheet(args.path, args.sheet, args.head)
    print(f"=== {args.path} / 工作表: {sheet}（{len(rows)} 行）===")
    for r in rows:
        print("| " + " | ".join(r) + " |")


def cmd_md2xlsx(args):
    with open(args.md, encoding="utf-8") as f:
        tables = extract_tables(f.read())
    if not tables:
        print("未找到 markdown 表格"); return
    sel = list(range(len(tables))) if args.table == "all" else [int(t) for t in args.table.split(",")]
    wb = Workbook()
    wb.remove(wb.active)
    for i in sel:
        rows = md_table_to_rows(tables[i])
        title = f"表{i+1}" if len(sel) > 1 else "Sheet1"
        ws = wb.create_sheet(title=title)
        for r in rows:
            ws.append(r)
    wb.save(args.out)
    print(f"已生成 {args.out}（{len(sel)} 张表）")


def cmd_xlsx2md(args):
    sheet, rows = read_sheet(args.path, args.sheet)
    if not rows:
        print("空表"); return
    ncol = max(len(r) for r in rows)
    print("| " + " | ".join(rows[0]) + " |")
    print("|" + "---|" * ncol)
    for r in rows[1:]:
        print("| " + " | ".join(r) + " |")


def cmd_stats(args):
    sheet, rows = read_sheet(args.path, args.sheet)
    if len(rows) < 2:
        print("空表或无数据行"); return
    header = rows[0]
    cols = args.col.split(",") if args.col else None
    print(f"=== 数值列统计: {args.path} / {sheet} ===")
    for i, h in enumerate(header):
        if cols and h not in cols:
            continue
        vals = []
        for r in rows[1:]:
            try:
                if i < len(r) and r[i].strip():
                    vals.append(float(r[i]))
            except ValueError:
                pass
        if not vals:
            continue
        n = len(vals)
        mean = sum(vals) / n
        srt = sorted(vals)
        med = srt[n // 2] if n % 2 else (srt[n // 2 - 1] + srt[n // 2]) / 2
        print(f"{h}: n={n} mean={mean:.3f} median={med:.3f} min={min(vals)} max={max(vals)}")


def cmd_csv2xlsx(args):
    wb = Workbook()
    ws = wb.active
    with open(args.csv, encoding="utf-8", newline="") as f:
        for row in csv.reader(f):
            ws.append(row)
    wb.save(args.out)
    print(f"已生成 {args.out}")


# ---------- markdown → Word/PPT（pandoc） ----------

PANDOC_FALLBACKS = [
    os.path.expanduser(r"~\AppData\Local\Pandoc\pandoc.exe"),  # winget 用户级安装
    r"C:\Program Files\Pandoc\pandoc.exe",
    r"C:\Program Files (x86)\Pandoc\pandoc.exe",
    "/usr/local/bin/pandoc",
    "/usr/bin/pandoc",
]


def find_pandoc():
    """PATH 优先；装完 winget 后 shell PATH 可能未刷新，兜底已知路径。"""
    exe = shutil.which("pandoc")
    if exe:
        return exe
    for p in PANDOC_FALLBACKS:
        if os.path.exists(p):
            return p
    return None


def run_pandoc(md, out, extra=None):
    pd = find_pandoc()
    if not pd:
        sys.exit("未找到 pandoc：请先安装（winget install JohnMacFarlane.Pandoc）")
    cmd = [pd, md, "-o", out] + (extra or [])
    r = subprocess.run(cmd, capture_output=True, text=True, encoding="utf-8")
    if r.returncode != 0:
        sys.exit(f"pandoc 失败:\n{r.stderr.strip()}")
    print(f"已生成 {out}")


def cmd_md2docx(args):
    extra = []
    if args.toc:
        extra += ["--toc", "--toc-depth=2"]
    if args.reference_doc:
        extra += [f"--reference-doc={args.reference_doc}"]
    run_pandoc(args.md, args.out, extra)


def cmd_md2pptx(args):
    extra = [f"--slide-level={args.slide_level}"]
    if args.reference_doc:
        extra += [f"--reference-doc={args.reference_doc}"]
    run_pandoc(args.md, args.out, extra)


def cmd_extract(args):
    os.makedirs(args.outdir, exist_ok=True)
    if args.fmt == "pdf":
        extract_pdf(args.path, args.outdir, args.pages)
    elif args.fmt == "docx":
        extract_docx(args.path, args.outdir)
    else:
        extract_pptx(args.path, args.outdir)


def extract_pdf(path, outdir, pages=None, min_size=250, min_kb=5, dpi=200):
    """提取 PDF 图 + 免费过滤。按页面区域渲染（贴合论文原图，含矢量/文字叠加），避免 raw 字节碎片。"""
    import fitz
    doc = fitz.open(path)
    kept = skipped = 0
    if pages:  # 渲染整页（用户指定）
        for p in pages.split(","):
            idx = int(p.strip())
            pix = doc[idx].get_pixmap(dpi=150)
            fn = os.path.join(outdir, f"page_{idx:03d}.png")
            pix.save(fn)
            print(f"渲染页 {idx} → {fn}")
            kept += 1
        print(f"\n渲染 {kept} 页完成")
        return
    for pno in range(len(doc)):
        page = doc[pno]
        ph = page.rect.height
        caps = [fitz.Rect(b[:4]) for b in page.get_text("blocks")
                if re.search(r"\b(Figure|Fig\.?)\s*\d+|图\s*\d+", b[4])]
        groups = []  # [union_rect, [(xref, info, rect)]]
        for img in page.get_images(full=True):
            xref = img[0]
            info = doc.extract_image(xref)
            w, h = info["width"], info["height"]
            if max(w, h) < min_size or len(info["image"]) < min_kb * 1024:
                skipped += 1
                continue  # 小图标/照片/薄装饰条
            rects = page.get_image_rects(xref)
            if not rects:
                skipped += 1
                continue
            r = rects[0]
            if r.y0 < ph * 0.04 or r.y1 > ph * 0.96:
                skipped += 1
                continue  # 页眉页脚带
            # 合并同面板相近图（相交 或 同行近距）
            for g in groups:
                if g[0].intersects(r) or (abs(g[0].y0 - r.y0) < 20 and abs(g[0].x0 - r.x0) < 60):
                    g[0] |= r
                    g[1].append((xref, info, r))
                    break
            else:
                groups.append([fitz.Rect(r), [(xref, info, r)]])
        for gi, (union, items) in enumerate(groups):
            m = 30  # 边距，尽量含图内文字/代码
            clip = fitz.Rect(max(0, union.x0 - m), max(0, union.y0 - m),
                             min(page.rect.width, union.x1 + m), min(page.rect.height, union.y1 + m))
            near_cap = any(clip.intersects(c) or abs(c.y0 - clip.y1) < 50 or abs(c.y1 - clip.y0) < 50 for c in caps)
            pix = page.get_pixmap(clip=clip, dpi=dpi)
            fn = os.path.join(outdir, f"p{pno}_fig{gi}.png")
            pix.save(fn)
            print(f"[{'CAP' if near_cap else '   '}] {fn} ({pix.width}x{pix.height}, {len(items)}图合并)")
            kept += 1
    print(f"\n保留 {kept} 组，跳过 {skipped} 张；按页面区域渲染，贴合论文原图（含矢量/文字叠加）")


def extract_docx(path, outdir):
    from docx import Document
    doc = Document(path)
    for i, rel in enumerate(doc.part.rels.values()):
        if "image" in rel.reltype:
            ext = rel.target_part.content_type.split("/")[-1].replace("jpeg", "jpg")
            fn = os.path.join(outdir, f"docx_img{i}.{ext}")
            with open(fn, "wb") as f:
                f.write(rel.target_part.blob)
            print(f"提取 {fn}")


def extract_pptx(path, outdir):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(path)
    n = 0
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                fn = os.path.join(outdir, f"slide{si}_img{n}.{shape.image.ext}")
                with open(fn, "wb") as f:
                    f.write(shape.image.blob)
                n += 1
                print(f"提取 {fn}")


def cmd_classify(args):
    """对幸存图用 vision 分类：价值图 vs 装饰图（只在 extract 过滤后跑，省 vision 调用）。"""
    import glob, subprocess
    files = sorted(glob.glob(os.path.join(args.dir, "*.png")) + glob.glob(os.path.join(args.dir, "*.jpg"))
                   + glob.glob(os.path.join(args.dir, "*.jpeg")))
    vision = args.vision or os.path.join(os.path.dirname(os.path.abspath(__file__)), "vision.py")
    if not files:
        print("目录里没有图"); return
    print(f"=== vision 分类 {len(files)} 张幸存图 ===")
    prompt = ("一句话判断：这张图是信息图（架构图/数据图/机制图/表格，有内容价值）"
              "还是装饰图（logo/照片/美化图标，无内容价值）？"
              "输出格式：价值图 或 装饰图，然后一句话理由。")
    for f in files:
        r = subprocess.run([sys.executable, vision, f, prompt],
                           capture_output=True, text=True, encoding="utf-8")
        out = r.stdout.strip() or r.stderr.strip()
        print(f"--- {os.path.basename(f)} ---")
        print(out)


def main():
    p = argparse.ArgumentParser(description="Office 文件处理（Excel 优先）")
    sub = p.add_subparsers(dest="cmd", required=True)
    r = sub.add_parser("read", help="读 xlsx 工作表")
    r.add_argument("path"); r.add_argument("--sheet"); r.add_argument("--head", type=int)
    m = sub.add_parser("md2xlsx", help="markdown 表 → xlsx")
    m.add_argument("md"); m.add_argument("out"); m.add_argument("--table", default="all", help="all 或 0,2 索引")
    x = sub.add_parser("xlsx2md", help="xlsx → markdown 表")
    x.add_argument("path"); x.add_argument("--sheet")
    s = sub.add_parser("stats", help="数值列统计")
    s.add_argument("path"); s.add_argument("--sheet"); s.add_argument("--col", help="逗号分隔列名")
    c = sub.add_parser("csv2xlsx", help="CSV → xlsx")
    c.add_argument("csv"); c.add_argument("out")
    e = sub.add_parser("extract", help="提取 PDF/docx/pptx 里的图（免费过滤装饰图）")
    e.add_argument("fmt", choices=["pdf", "docx", "pptx"])
    e.add_argument("path")
    e.add_argument("--outdir", default="figures")
    e.add_argument("--pages", help="pdf 渲染指定页(逗号分隔)，默认提取嵌入图")
    e.add_argument("--min-size", type=int, default=250, help="最小边长px，滤小图")
    e.add_argument("--min-kb", type=int, default=5, help="最小文件KB，滤装饰条")
    cl = sub.add_parser("classify", help="vision 分类幸存图：价值 vs 装饰")
    cl.add_argument("dir")
    cl.add_argument("--vision", help="vision.py 路径（默认同目录）")
    d = sub.add_parser("md2docx", help="markdown → Word（pandoc；公式转 OMML 原生方程）")
    d.add_argument("md"); d.add_argument("out")
    d.add_argument("--toc", action="store_true", help="加目录域（Word 里打开更新）")
    d.add_argument("--reference-doc", help="样式模板 docx（--reference-doc 控中文字体/缩进）")
    k = sub.add_parser("md2pptx", help="markdown → PPT（pandoc；文本框可编辑）")
    k.add_argument("md"); k.add_argument("out")
    k.add_argument("--slide-level", type=int, default=2, help="几级标题=一页（默认2：#分节、##一页）")
    k.add_argument("--reference-doc", help="样式模板 pptx")
    args = p.parse_args()
    {"read": cmd_read, "md2xlsx": cmd_md2xlsx, "xlsx2md": cmd_xlsx2md,
     "stats": cmd_stats, "csv2xlsx": cmd_csv2xlsx, "extract": cmd_extract,
     "classify": cmd_classify, "md2docx": cmd_md2docx, "md2pptx": cmd_md2pptx}[args.cmd](args)


if __name__ == "__main__":
    main()
